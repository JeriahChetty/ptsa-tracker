"""
Document parser for extracting measure data from PDF and PowerPoint files
Supports both pattern matching and AI-powered extraction
"""
import os
import re
import base64
from datetime import datetime
from typing import Dict, List, Optional
from PyPDF2 import PdfReader
from pptx import Presentation
from io import BytesIO


def parse_measure_document(file_path: str, file_type: str, use_ai: bool = True) -> Dict:
    """
    Parse a PDF, PowerPoint, or image file and extract measure information
    
    Args:
        file_path: Path to the uploaded file
        file_type: 'pdf', 'pptx', 'png', 'jpg', 'jpeg', 'webp'
        use_ai: Whether to use AI-powered extraction (requires OPENAI_API_KEY)
    
    Returns:
        Dictionary with extracted measure data (can contain multiple measures)
    """
    # Handle image files
    if file_type in ['png', 'jpg', 'jpeg', 'webp']:
        if use_ai and os.getenv('OPENAI_API_KEY'):
            try:
                measures = parse_image_with_vision(file_path)
                if measures:
                    return {'measures': measures, 'method': 'ai_vision'}
            except Exception as e:
                print(f"AI Vision extraction failed: {e}")
                return {'measures': [], 'method': 'error', 'error': str(e)}
        else:
            return {'measures': [], 'method': 'error', 'error': 'OpenAI API key required for image parsing'}
    
    # Handle PDF/PowerPoint
    if file_type == 'pdf':
        text, images = parse_pdf(file_path)
    elif file_type in ['pptx', 'ppt']:
        text, images = parse_powerpoint(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    # Try AI extraction first if enabled and API key is available
    if use_ai and os.getenv('OPENAI_API_KEY'):
        try:
            measures = extract_with_ai(text, images)
            if measures:
                return {'measures': measures, 'method': 'ai'}
        except Exception as e:
            print(f"AI extraction failed, falling back to pattern matching: {e}")
    
    # Fallback to pattern matching
    measures = extract_multiple_measures(text)
    return {'measures': measures, 'method': 'pattern_matching'}


def parse_pdf(file_path: str) -> tuple:
    """Extract text and images from PDF"""
    reader = PdfReader(file_path)
    text = ""
    images = []
    
    # Extract text from all pages
    for page in reader.pages:
        text += page.extract_text() + "\n"
        
        # Extract images from page
        try:
            if '/XObject' in page['/Resources']:
                xObject = page['/Resources']['/XObject'].get_object()
                for obj in xObject:
                    if xObject[obj]['/Subtype'] == '/Image':
                        try:
                            size = (xObject[obj]['/Width'], xObject[obj]['/Height'])
                            data = xObject[obj].get_data()
                            images.append({'data': data, 'size': size})
                        except:
                            pass
        except:
            pass
    
    return text, images


def parse_powerpoint(file_path: str) -> tuple:
    """Extract text and images from PowerPoint"""
    prs = Presentation(file_path)
    text = ""
    images = []
    
    # Extract text and images from all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            # Extract images
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    images.append({'data': image_bytes, 'size': None})
                except:
                    pass
    
    return text, images


def parse_image_with_vision(file_path: str) -> List[Dict]:
    """
    Parse an image file using OpenAI Vision API to extract measure data
    
    Args:
        file_path: Path to the image file
    
    Returns:
        List of measures extracted from the image
    """
    from openai import OpenAI
    
    client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    # Read and encode image
    with open(file_path, 'rb') as image_file:
        image_data = base64.b64encode(image_file.read()).decode('utf-8')
    
    # Determine image format from file extension
    file_ext = file_path.split('.')[-1].lower()
    mime_type = f"image/{file_ext if file_ext != 'jpg' else 'jpeg'}"
    
    prompt = """
    Analyze this image of a business improvement measure document and extract ALL measures.
    Return a JSON object with a "measures" array, where each measure has these fields:
    
    - name: Measure name/title
    - measure_detail: Description of the measure
    - target: Target or objective
    - departments: Focus area or departments (e.g., "Process", "Technology")
    - responsible: Person responsible
    - participants: Comma-separated list of participants
    - start_date: Start date in YYYY-MM-DD format (if mentioned)
    - end_date: End date in YYYY-MM-DD format (if mentioned)
    - steps: Array of step descriptions (extract "Step 1", "Step 2", etc.)
    
    Parse all text visible in the image, including tables. If a field is not found, use null.
    Extract dates intelligently (e.g., "October 30, 2025" -> "2025-10-30").
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_data}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=2000,
            temperature=0.1
        )
        
        import json
        result = json.loads(response.choices[0].message.content)
        
        # Handle different response formats
        if 'measures' in result:
            measures = result['measures']
        elif isinstance(result, list):
            measures = result
        else:
            measures = [result]
        
        # Convert date strings to date objects
        for measure in measures:
            if measure.get('start_date'):
                try:
                    measure['start_date'] = datetime.fromisoformat(measure['start_date']).date()
                except:
                    measure['start_date'] = parse_date(measure['start_date'])
            
            if measure.get('end_date'):
                try:
                    measure['end_date'] = datetime.fromisoformat(measure['end_date']).date()
                except:
                    measure['end_date'] = parse_date(measure['end_date'])
        
        return measures
        
    except Exception as e:
        print(f"OpenAI Vision API error: {e}")
        raise


def extract_with_ai(text: str, images: List[Dict]) -> List[Dict]:
    """
    Use OpenAI API to extract measure data from text
    Returns a list of measures (supports multiple measures in one document)
    """
    from openai import OpenAI
    
    client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    prompt = """
    Extract all improvement measures from the following document text. 
    Return a JSON array of measure objects, where each measure has these fields:
    
    - name: Measure name/title
    - measure_detail: Description of the measure
    - target: Target or objective
    - departments: Focus area or departments involved
    - responsible: Person responsible
    - participants: Comma-separated list of participants
    - start_date: Start date in YYYY-MM-DD format (if mentioned)
    - end_date: End date in YYYY-MM-DD format (if mentioned)
    - steps: Array of step descriptions (Step 1, Step 2, etc.)
    
    If a field is not found, use null. Parse dates intelligently.
    
    Document text:
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4-turbo-preview",
            messages=[
                {"role": "system", "content": "You are an expert at extracting structured data from business improvement documents. Always return valid JSON."},
                {"role": "user", "content": prompt + "\n\n" + text}
            ],
            response_format={"type": "json_object"},
            temperature=0.1
        )
        
        import json
        result = json.loads(response.choices[0].message.content)
        
        # Handle different possible response formats
        if 'measures' in result:
            measures = result['measures']
        elif isinstance(result, list):
            measures = result
        else:
            # Assume it's a single measure
            measures = [result]
        
        # Convert date strings to date objects
        for measure in measures:
            if measure.get('start_date'):
                try:
                    measure['start_date'] = datetime.fromisoformat(measure['start_date']).date()
                except:
                    measure['start_date'] = parse_date(measure['start_date'])
            
            if measure.get('end_date'):
                try:
                    measure['end_date'] = datetime.fromisoformat(measure['end_date']).date()
                except:
                    measure['end_date'] = parse_date(measure['end_date'])
        
        return measures
        
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return []


def extract_multiple_measures(text: str) -> List[Dict]:
    """
    Extract multiple measures from text using pattern matching
    Looks for "Measure 1", "Measure 2", etc. to split the document
    """
    measures = []
    
    # Try to split by measure numbers
    measure_pattern = r'Measure\s+(\d+)[:\s]+'
    splits = list(re.finditer(measure_pattern, text, re.IGNORECASE))
    
    if len(splits) > 1:
        # Multiple measures found
        for i, match in enumerate(splits):
            start = match.start()
            end = splits[i + 1].start() if i + 1 < len(splits) else len(text)
            measure_text = text[start:end]
            measure_data = extract_measure_data(measure_text)
            if measure_data.get('name'):
                measures.append(measure_data)
    else:
        # Single measure or no clear separation
        measure_data = extract_measure_data(text)
        if measure_data.get('name'):
            measures.append(measure_data)
    
    return measures


def extract_measure_data(text: str) -> Dict:
    """
    Extract structured measure data from text using pattern matching
    
    Expected patterns based on the image:
    - Measure name/title at the top
    - Focus area: Process/Technology/etc
    - Description: [text]
    - Responsible: [name]
    - Participants: [names]
    - Time/Date: [date]
    - Target: [text]
    - Steps: Step 1, Step 2, etc.
    """
    data = {
        'name': '',
        'measure_detail': '',
        'target': '',
        'departments': '',
        'responsible': '',
        'participants': '',
        'start_date': None,
        'end_date': None,
        'steps': []
    }
    
    # Extract measure name (usually first line or after "Measure")
    measure_match = re.search(r'Measure\s+\d+\s*[:\-]?\s*(.+?)(?:\n|$)', text, re.IGNORECASE)
    if measure_match:
        data['name'] = measure_match.group(1).strip()
    else:
        # Try to get first non-empty line
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        if lines:
            data['name'] = lines[0]
    
    # Extract focus area / departments
    focus_match = re.search(r'Focus\s+area\s*:?\s*(.+?)(?:\n|Description)', text, re.IGNORECASE)
    if focus_match:
        data['departments'] = focus_match.group(1).strip()
    
    # Extract description
    desc_match = re.search(r'Description\s*:?\s*(.+?)(?:\n(?:Responsible|Participants|Time|Target|Step)|$)', text, re.IGNORECASE | re.DOTALL)
    if desc_match:
        data['measure_detail'] = desc_match.group(1).strip()
    
    # Extract responsible person
    resp_match = re.search(r'Responsible\s*:?\s*(.+?)(?:\n|$)', text, re.IGNORECASE)
    if resp_match:
        data['responsible'] = resp_match.group(1).strip()
    
    # Extract participants
    part_match = re.search(r'Participants\s*:?\s*(.+?)(?:\n|$)', text, re.IGNORECASE)
    if part_match:
        data['participants'] = part_match.group(1).strip()
    
    # Extract target
    target_match = re.search(r'Target\s*:?\s*(.+?)(?:\n|Step)', text, re.IGNORECASE | re.DOTALL)
    if target_match:
        data['target'] = target_match.group(1).strip()
    
    # Extract date/time
    time_match = re.search(r'Time\s*:?\s*(.+?)(?:\n|$)', text, re.IGNORECASE)
    if time_match:
        date_str = time_match.group(1).strip()
        # Try to parse date
        parsed_date = parse_date(date_str)
        if parsed_date:
            data['end_date'] = parsed_date
    
    # Extract steps
    steps = extract_steps(text)
    data['steps'] = steps
    
    return data


def extract_steps(text: str) -> List[str]:
    """Extract step titles from text"""
    steps = []
    
    # Pattern: Step 1, Step 2, etc.
    step_pattern = r'Step\s+\d+\s*:?\s*(.+?)(?=\nStep\s+\d+|\n\n|$)'
    matches = re.finditer(step_pattern, text, re.IGNORECASE | re.DOTALL)
    
    for match in matches:
        step_text = match.group(1).strip()
        # Clean up step text (remove extra whitespace, newlines)
        step_text = ' '.join(step_text.split())
        if step_text:
            steps.append(step_text)
    
    return steps


def parse_date(date_str: str) -> Optional[datetime]:
    """
    Try to parse a date string in various formats
    
    Examples:
    - October 30, 2025
    - 2025-10-30
    - 30/10/2025
    """
    date_formats = [
        '%B %d, %Y',      # October 30, 2025
        '%b %d, %Y',      # Oct 30, 2025
        '%Y-%m-%d',       # 2025-10-30
        '%d/%m/%Y',       # 30/10/2025
        '%m/%d/%Y',       # 10/30/2025
        '%d-%m-%Y',       # 30-10-2025
        '%Y/%m/%d',       # 2025/10/30
    ]
    
    # Extract just the date part if there's extra text
    date_match = re.search(r'([A-Za-z]+\s+\d{1,2},?\s+\d{4}|\d{1,2}[-/]\d{1,2}[-/]\d{4}|\d{4}[-/]\d{1,2}[-/]\d{1,2})', date_str)
    if date_match:
        date_str = date_match.group(1)
    
    for fmt in date_formats:
        try:
            return datetime.strptime(date_str.strip(), fmt).date()
        except ValueError:
            continue
    
    return None
